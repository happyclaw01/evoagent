# Copyright (c) 2025 MiroMind
# This source code is licensed under the MIT License.

import html
import json
import os
import re
import shutil
import subprocess
import tempfile
import traceback
from typing import Any, Dict, List, Union
from urllib.parse import parse_qs, quote, unquote, urlparse, urlunparse

import mammoth
import markdownify
import openpyxl
import pdfminer
import pdfminer.high_level
import pptx
from bs4 import BeautifulSoup
from markitdown import MarkItDown
from openpyxl.utils import get_column_letter
from youtube_transcript_api._api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import SRTFormatter


def process_input(task_description, task_file_name):
    """
    Process user input, especially files.
    Returns formatted initial user message content list and updated task description.
    """
    initial_user_content = ""
    updated_task_description = task_description

    if task_file_name:
        try:
            file_extension = task_file_name.rsplit(".", maxsplit=1)[-1].lower()
            parsing_result = None

            if file_extension in ["jpg", "jpeg", "png", "gif", "webp"]:
                updated_task_description += f"\nNote: An Image file '{task_file_name}' is associated with this task. You may use available tools to read its content if necessary.\n\n"

            elif file_extension == "txt":
                with open(task_file_name, "r") as f:
                    initial_user_content += f.read()
                updated_task_description += f"\nNote: An Excel file '{task_file_name}' is associated with this task. You may use available tools to read its content if necessary.\n\n"

            elif file_extension in ["jsonld", "json"]:
                with open(task_file_name, "r") as f:
                    initial_user_content += json.dumps(json.load(f), ensure_ascii=False)
                updated_task_description += f"\nNote: An Json file '{task_file_name}' is associated with this task. You may use available tools to read its content if necessary.\n\n"

            elif file_extension in ["xlsx", "xls"]:
                parsing_result = XlsxConverter(local_path=task_file_name)
                # Add note for Excel files to inform LLM that tools can be used to read
                updated_task_description += f"\nNote: An Excel file '{task_file_name}' is associated with this task. You may use available tools to read its content if necessary.\n\n"

            elif file_extension == "pdf":
                parsing_result = DocumentConverterResult(
                    title=None,
                    text_content=pdfminer.high_level.extract_text(task_file_name),
                )
                updated_task_description += f"\nNote: A PDF file '{task_file_name}' is associated with this task. You may use available tools to read its content if necessary.\n\n"

            elif file_extension in ["docx", "doc"]:
                parsing_result = DocxConverter(local_path=task_file_name)
                updated_task_description += f"\nNote: A Document file '{task_file_name}' is associated with this task. You may use available tools to read its content if necessary.\n\n"

            elif file_extension in ["html", "htm"]:
                parsing_result = HtmlConverter(local_path=task_file_name)
                updated_task_description += f"\nNote: An HTML file '{task_file_name}' is associated with this task. You may use available tools to read its content if necessary.\n\n"

            elif file_extension in ["pptx", "ppt"]:
                parsing_result = PptxConverter(local_path=task_file_name)
                updated_task_description += f"\nNote: A PPT file '{task_file_name}' is associated with this task. You may use available tools to read its content if necessary.\n\n"

            elif file_extension in ["wav"]:
                updated_task_description += f"\nNote: An audio file '{task_file_name}' is associated with this task. You may use available tools to read its content if necessary.\n\n"

            elif file_extension in ["mp3", "m4a"]:
                updated_task_description += f"\nNote: An audio file '{task_file_name}' is associated with this task. You may use available tools to read its content if necessary.\n\n"

            elif file_extension in ["zip"]:
                parsing_result = ZipConverter(local_path=task_file_name)
                updated_task_description += f"\nNote: Several files zipped in '{task_file_name}' are associated with this task. You may use available tools to read its content if necessary.\n\n"

            else:
                # For other file types, consider adding general prompts or specific processing
                updated_task_description += f"\nNote: A file '{task_file_name}' is associated with this task. You may use available tools to read its content if necessary.\n\n"

            #### markitdown process ####
            try:
                if file_extension not in [
                    "jpg",
                    "jpeg",
                    "png",
                    "gif",
                    "webp",
                    "wav",
                    "mp3",
                    "pdb",
                ]:
                    from markitdown import MarkItDown

                    md = MarkItDown(enable_plugins=True)
                    parsing_result = md.convert(task_file_name)
                    print(f"Info: Used MarkItDown to process file {task_file_name}")
            except Exception:
                pass
            ############################

            # add the content and title (if has) into the task description
            if getattr(parsing_result, "title", None):
                updated_task_description += (
                    "<file>Title:\n{}\n\nContent:\n{}</file>".format(
                        parsing_result.title, parsing_result.text_content
                    )
                )
            elif getattr(parsing_result, "text_content", None):
                content = parsing_result.text_content
                max_len = 200_000  # Limit the length of results returned to LLM
                if len(content) > max_len:
                    content = content[:max_len] + "\n... [File truncated]"
                updated_task_description += "<file>{}</file>".format(content)
            else:
                pass  # for image file

        except FileNotFoundError:
            raise (f"Error: File not found {task_file_name}")
            updated_task_description += (
                f"\nWarning: The specified file '{task_file_name}' was not found."
            )
        except Exception as e:
            raise (f"Error: Error processing file {task_file_name}: {e}")
            updated_task_description += (
                f"\nWarning: There was an error processing the file '{task_file_name}'."
            )

    # output format requirement
    use_cn_prompt = os.getenv("USE_CN_PROMPT", "0")
    if use_cn_prompt == "1":
        updated_task_description += "\n请通过任务分解和MCP工具调用来解决给定的问题。**你必须严格遵循请求中的格式要求，并将最终答案包裹在 \\boxed{} 中。**"
    else:
        updated_task_description += "\nYou should follow the format instruction in the request strictly and wrap the final answer in \\boxed{}."
    initial_user_content += updated_task_description

    return initial_user_content, updated_task_description


class _CustomMarkdownify(markdownify.MarkdownConverter):
    """
    A custom version of markdownify's MarkdownConverter. Changes include:

    - Altering the default heading style to use '#', '##', etc.
    - Removing javascript hyperlinks.
    - Truncating images with large data:uri sources.
    - Ensuring URIs are properly escaped, and do not conflict with Markdown syntax
    """

    def __init__(self, **options: Any):
        options["heading_style"] = options.get("heading_style", markdownify.ATX)
        # Explicitly cast options to the expected type if necessary
        super().__init__(**options)

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual, but be sure to start with a new line"""
        if not convert_as_inline:
            if not re.search(r"^\n", text):
                return "\n" + super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

        return super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

    def convert_a(self, el: Any, text: str, convert_as_inline: bool):
        """Same as usual converter, but removes Javascript links and escapes URIs."""
        prefix, suffix, text = markdownify.chomp(text)  # type: ignore
        if not text:
            return ""
        href = el.get("href")
        title = el.get("title")

        # Escape URIs and skip non-http or file schemes
        if href:
            try:
                parsed_url = urlparse(href)  # type: ignore
                if parsed_url.scheme and parsed_url.scheme.lower() not in [
                    "http",
                    "https",
                    "file",
                ]:  # type: ignore
                    return "%s%s%s" % (prefix, text, suffix)
                href = urlunparse(
                    parsed_url._replace(path=quote(unquote(parsed_url.path)))
                )  # type: ignore
            except ValueError:  # It's not clear if this ever gets thrown
                return "%s%s%s" % (prefix, text, suffix)

        # For the replacement see #29: text nodes underscores are escaped
        if (
            self.options["autolinks"]
            and text.replace(r"\_", "_") == href
            and not title
            and not self.options["default_title"]
        ):
            # Shortcut syntax
            return "<%s>" % href
        if self.options["default_title"] and not title:
            title = href
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        return (
            "%s[%s](%s%s)%s" % (prefix, text, href, title_part, suffix)
            if href
            else text
        )

    def convert_img(self, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual converter, but removes data URIs"""

        alt = el.attrs.get("alt", None) or ""
        src = el.attrs.get("src", None) or ""
        title = el.attrs.get("title", None) or ""
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        if (
            convert_as_inline
            and el.parent.name not in self.options["keep_inline_images_in"]
        ):
            return alt

        # Remove dataURIs
        if src.startswith("data:"):
            src = src.split(",")[0] + "..."

        return "![%s](%s%s)" % (alt, src, title_part)

    def convert_soup(self, soup: Any) -> str:
        return super().convert_soup(soup)  # type: ignore


class DocumentConverterResult:
    """The result of converting a document to text."""

    def __init__(self, title: Union[str, None] = None, text_content: str = ""):
        self.title: Union[str, None] = title
        self.text_content: str = text_content


def convert_html_to_md(html_content):
    """
    Placeholder for HTML to Markdown conversion function
    In the original class, this would call self._convert()
    """
    soup = BeautifulSoup(html_content, "html.parser")
    for script in soup(["script", "style"]):
        script.extract()

    # Print only the main content
    body_elm = soup.find("body")
    webpage_text = ""
    if body_elm:
        webpage_text = _CustomMarkdownify().convert_soup(body_elm)
    else:
        webpage_text = _CustomMarkdownify().convert_soup(soup)

    assert isinstance(webpage_text, str)

    return DocumentConverterResult(
        title=None if soup.title is None else soup.title.string,
        text_content=webpage_text,
    )


def HtmlConverter(local_path: str):
    with open(local_path, "rt", encoding="utf-8") as fh:
        html_content = fh.read()

        return convert_html_to_md(html_content)


def DocxConverter(local_path: str):
    with open(local_path, "rb") as docx_file:
        result = mammoth.convert_to_html(docx_file)
        html_content = result.value
    return convert_html_to_md(html_content)


def WikipediaConverter(local_path: str, **kwargs):
    url = kwargs.get("url", "")
    assert re.search(r"^https?:\/\/[a-zA-Z]{2,3}\.wikipedia.org\/", url) is not None

    # Parse the file
    soup = None
    with open(local_path, "rt", encoding="utf-8") as fh:
        soup = BeautifulSoup(fh.read(), "html.parser")

    # Remove javascript and style blocks
    for script in soup(["script", "style"]):
        script.extract()

    # Print only the main content
    body_elm = soup.find("div", {"id": "mw-content-text"})
    title_elm = soup.find("span", {"class": "mw-page-title-main"})

    webpage_text = ""
    main_title = None if soup.title is None else soup.title.string

    if body_elm:
        # What's the title
        if title_elm and len(title_elm) > 0:
            main_title = title_elm.string  # type: ignore
            assert isinstance(main_title, str)

        # Convert the page
        webpage_text = f"# {main_title}\n\n" + _CustomMarkdownify().convert_soup(
            body_elm
        )
    else:
        webpage_text = _CustomMarkdownify().convert_soup(soup)

    return DocumentConverterResult(
        title=main_title,
        text_content=webpage_text,
    )


def YouTubeConverter(local_path: str, url: str):
    assert url.startswith("https://www.youtube.com/watch?")

    def _get(
        metadata: Dict[str, str], keys: List[str], default: Union[str, None] = None
    ) -> Union[str, None]:
        for k in keys:
            if k in metadata:
                return metadata[k]
        return default

    def _findKey(json: Any, key: str) -> Union[str, None]:  # TODO: Fix json type
        if isinstance(json, list):
            for elm in json:
                ret = _findKey(elm, key)
                if ret is not None:
                    return ret
        elif isinstance(json, dict):
            for k in json:
                if k == key:
                    return json[k]
                else:
                    ret = _findKey(json[k], key)
                    if ret is not None:
                        return ret
        return None

    # Parse the file
    soup = None
    with open(local_path, "rt", encoding="utf-8") as fh:
        soup = BeautifulSoup(fh.read(), "html.parser")

    # Read the meta tags
    assert soup.title is not None and soup.title.string is not None
    metadata: Dict[str, str] = {"title": soup.title.string}
    for meta in soup(["meta"]):
        for a in meta.attrs:
            if a in ["itemprop", "property", "name"]:
                metadata[meta[a]] = meta.get("content", "")
                break

    # We can also try to read the full description. This is more prone to breaking, since it reaches into the page implementation
    try:
        for script in soup(["script"]):
            content = script.text
            if "ytInitialData" in content:
                lines = re.split(r"\r?\n", content)
                obj_start = lines[0].find("{")
                obj_end = lines[0].rfind("}")
                if obj_start >= 0 and obj_end >= 0:
                    data = json.loads(lines[0][obj_start : obj_end + 1])
                    attrdesc = _findKey(data, "attributedDescriptionBodyText")  # type: ignore
                    if attrdesc:
                        metadata["description"] = str(attrdesc["content"])
                break
    except Exception:
        pass

    # Start preparing the page
    webpage_text = "# YouTube\n"

    title = _get(metadata, ["title", "og:title", "name"])  # type: ignore
    assert isinstance(title, str)

    if title:
        webpage_text += f"\n## {title}\n"

    stats = ""
    views = _get(metadata, ["interactionCount"])  # type: ignore
    if views:
        stats += f"- **Views:** {views}\n"

    keywords = _get(metadata, ["keywords"])  # type: ignore
    if keywords:
        stats += f"- **Keywords:** {keywords}\n"

    runtime = _get(metadata, ["duration"])  # type: ignore
    if runtime:
        stats += f"- **Runtime:** {runtime}\n"

    if len(stats) > 0:
        webpage_text += f"\n### Video Metadata\n{stats}\n"

    description = _get(metadata, ["description", "og:description"])  # type: ignore
    if description:
        webpage_text += f"\n### Description\n{description}\n"

    transcript_text = ""
    parsed_url = urlparse(url)  # type: ignore
    params = parse_qs(parsed_url.query)  # type: ignore
    if "v" in params:
        assert isinstance(params["v"][0], str)
        video_id = str(params["v"][0])
        try:
            # Must be a single transcript.
            transcript = YouTubeTranscriptApi.get_transcript(video_id)  # type: ignore
            # transcript_text = " ".join([part["text"] for part in transcript])  # type: ignore
            # Alternative formatting:
            transcript_text = SRTFormatter().format_transcript(transcript)
        except Exception:
            pass
    if transcript_text:
        webpage_text += f"\n### Transcript\n{transcript_text}\n"

    title = title if title else soup.title.string
    assert isinstance(title, str)

    return DocumentConverterResult(
        title=title,
        text_content=webpage_text,
    )


def XlsxConverter(local_path: str):
    """
    Converts Excel files to Markdown using openpyxl.
    Preserves color formatting and other cell styling information.

    Args:
        local_path: Path to the Excel file

    Returns:
        DocumentConverterResult with the Markdown representation of the Excel file
    """
    # Load the workbook
    wb = openpyxl.load_workbook(local_path, data_only=True)
    md_content = ""

    # Helper function to convert RGB color to hex
    def rgb_to_hex(rgb_value):
        if not rgb_value:
            return None

        # Convert RGB value to string for processing
        rgb_string = str(rgb_value)

        # Handle RGB format like 'RGB(255, 255, 255)'
        if isinstance(rgb_value, str) and rgb_string.startswith("RGB"):
            rgb_match = re.match(r"RGB\((\d+), (\d+), (\d+)\)", rgb_string)
            if rgb_match:
                r, g, b = map(int, rgb_match.groups())
                return f"#{r:02x}{g:02x}{b:02x}"

        # Special handling for FFFFFFFF (white) and 00000000 (transparent/none)
        if rgb_string in ["FFFFFFFF", "00000000", "none", "auto"]:
            return None

        # Handle ARGB format (common in openpyxl)
        if len(rgb_string) == 8:  # ARGB format like 'FF5733FF'
            return f"#{rgb_string[2:]}"  # Strip alpha channel

        # Handle direct hex values like 'FF5733'
        if isinstance(rgb_value, str):
            return f"#{rgb_string}" if not rgb_string.startswith("#") else rgb_string

        return None  # Return None for unrecognized formats

    # Helper function to detect and format cell styling
    def get_cell_format_info(cell):
        info = {}

        # Get background color if it exists
        if cell.fill and hasattr(cell.fill, "fgColor") and cell.fill.fgColor:
            # Get the RGB value - in openpyxl this can be stored in different attributes
            rgb_value = None
            if hasattr(cell.fill.fgColor, "rgb") and cell.fill.fgColor.rgb:
                rgb_value = cell.fill.fgColor.rgb
            elif hasattr(cell.fill.fgColor, "value") and cell.fill.fgColor.value:
                rgb_value = cell.fill.fgColor.value

            if rgb_value:
                bg_color = rgb_to_hex(rgb_value)
                if bg_color:  # Skip transparent or white (handled in rgb_to_hex)
                    info["bg_color"] = bg_color

        # Get font color if it exists
        if cell.font and hasattr(cell.font, "color") and cell.font.color:
            # Get the RGB value - in openpyxl this can be stored in different attributes
            rgb_value = None
            if hasattr(cell.font.color, "rgb") and cell.font.color.rgb:
                rgb_value = cell.font.color.rgb
            elif hasattr(cell.font.color, "value") and cell.font.color.value:
                rgb_value = cell.font.color.value

            if rgb_value:
                font_color = rgb_to_hex(rgb_value)
                if font_color:  # Skip transparent (handled in rgb_to_hex)
                    info["font_color"] = font_color

        # Get font weight (bold)
        if cell.font and cell.font.bold:
            info["bold"] = True

        # Get font style (italic)
        if cell.font and cell.font.italic:
            info["italic"] = True

        # Get font underline
        if cell.font and cell.font.underline and cell.font.underline != "none":
            info["underline"] = True

        return info

    # Process each sheet in the workbook
    for sheet_name in wb.sheetnames:
        try:
            sheet = wb[sheet_name]
            md_content += f"## {sheet_name}\n\n"

            # Get the dimensions of the used part of the sheet
            min_row, min_col = 1, 1
            max_row = max(
                (cell.row for cell in sheet._cells.values() if cell.value is not None),
                default=0,
            )
            max_col = max(
                (
                    cell.column
                    for cell in sheet._cells.values()
                    if cell.value is not None
                ),
                default=0,
            )

            if max_row == 0 or max_col == 0:
                md_content += "This sheet is empty.\n\n"
                continue
        except Exception as e:
            error_msg = f"Error processing sheet '{sheet_name}': {str(e)}"
            print(error_msg)
            md_content += (
                f"## {sheet_name}\n\nError processing this sheet: {str(e)}\n\n"
            )
            continue

        try:
            # First, determine column widths
            col_widths = {}
            for col_idx in range(min_col, max_col + 1):
                max_length = 0
                # col_letter = get_column_letter(col_idx)
                _ = get_column_letter(col_idx)
                for row_idx in range(min_row, max_row + 1):
                    try:
                        cell = sheet.cell(row=row_idx, column=col_idx)
                        cell_value = str(cell.value) if cell.value is not None else ""
                        max_length = max(max_length, len(cell_value))
                    except Exception as e:
                        print(
                            f"Warning: Error processing cell at row {row_idx}, column {col_idx}: {str(e)}"
                        )
                        max_length = max(max_length, 10)  # Use reasonable default
                col_widths[col_idx] = max(max_length + 2, 5)  # Min width of 5

            # Start building the table
            # Header row with column separators
            md_content += "|"
            for col_idx in range(min_col, max_col + 1):
                md_content += " " + " " * col_widths[col_idx] + " |"
            md_content += "\n"

            # Separator row
            md_content += "|"
            for col_idx in range(min_col, max_col + 1):
                md_content += ":" + "-" * col_widths[col_idx] + ":|"
            md_content += "\n"

            # Data rows
            for row_idx in range(min_row, max_row + 1):
                md_content += "|"
                for col_idx in range(min_col, max_col + 1):
                    try:
                        cell = sheet.cell(row=row_idx, column=col_idx)
                        cell_value = str(cell.value) if cell.value is not None else ""

                        # Get formatting info
                        try:
                            format_info = get_cell_format_info(cell)
                        except Exception as e:
                            print(
                                f"Warning: Error getting formatting for cell at row {row_idx}, column {col_idx}: {str(e)}"
                            )
                            format_info = {}

                        formatted_value = cell_value

                        # Add HTML-style formatting if needed
                        if format_info:
                            style_parts = []

                            if "bg_color" in format_info:
                                style_parts.append(
                                    f"background-color:{format_info['bg_color']}"
                                )

                            if "font_color" in format_info:
                                style_parts.append(f"color:{format_info['font_color']}")

                            span_attributes = []
                            if style_parts:
                                span_attributes.append(
                                    f'style="{"; ".join(style_parts)}"'
                                )

                            # Format with bold/italic/underline if needed
                            inner_value = cell_value
                            if "bold" in format_info:
                                inner_value = f"<strong>{inner_value}</strong>"
                            if "italic" in format_info:
                                inner_value = f"<em>{inner_value}</em>"
                            if "underline" in format_info:
                                inner_value = f"<u>{inner_value}</u>"

                            # Only add a span if we have style attributes
                            if span_attributes:
                                formatted_value = f"<span {' '.join(span_attributes)}>{inner_value}</span>"
                            else:
                                formatted_value = inner_value

                        # Pad to column width and add to markdown
                        padding = col_widths[col_idx] - len(cell_value)
                        padded_value = " " + formatted_value + " " * (padding + 1)
                        md_content += padded_value + "|"
                    except Exception as e:
                        print(
                            f"Error processing cell at row {row_idx}, column {col_idx}: {str(e)}"
                        )
                        # Add a placeholder for the failed cell
                        padded_value = " [Error] " + " " * (col_widths[col_idx] - 7)
                        md_content += padded_value + " |"

                md_content += "\n"
        except Exception as e:
            error_msg = f"Error generating table for sheet '{sheet_name}': {str(e)}\n{traceback.format_exc()}"
            print(error_msg)
            md_content += f"Error generating table: {str(e)}\n\n"

        # Add formatting legend
        has_formatting = False
        for row_idx in range(min_row, max_row + 1):
            for col_idx in range(min_col, max_col + 1):
                cell = sheet.cell(row=row_idx, column=col_idx)
                if get_cell_format_info(cell):
                    has_formatting = True
                    break
            if has_formatting:
                break

        if has_formatting:
            md_content += "\n### Formatting Information\n"
            md_content += "The table above includes HTML formatting to represent colors and styles from the original Excel file.\n"
            md_content += "This formatting may not display in all Markdown viewers.\n"

        md_content += "\n\n"  # Extra newlines between sheets

    return DocumentConverterResult(
        title=None,
        text_content=md_content.strip(),
    )


def PptxConverter(local_path) -> Union[None, dict]:
    """
    Converts PPTX files to Markdown. Supports headings, tables and images with alt text.

    Args:
        local_path: Path to the PPTX file
        file_extension: Extension of the file (default: ".pptx")

    Returns:
        None if not a PPTX file, otherwise a dictionary with title and text_content
    """

    def is_picture(shape):
        """Check if a shape is a picture"""
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def is_table(shape):
        """Check if a shape is a table"""
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    assert local_path.endswith(".pptx")

    md_content = ""
    presentation = pptx.Presentation(local_path)
    slide_num = 0

    for slide in presentation.slides:
        slide_num += 1
        md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"
        title = slide.shapes.title

        for shape in slide.shapes:
            # Pictures
            if is_picture(shape):
                # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069
                alt_text = ""
                try:
                    alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                except Exception:
                    pass
                # A placeholder name
                filename = re.sub(r"\W", "", shape.name) + ".jpg"
                md_content += (
                    "\n!["
                    + (alt_text if alt_text else shape.name)
                    + "]("
                    + filename
                    + ")\n"
                )

            # Tables
            if is_table(shape):
                html_table = "<html><body><table>"
                first_row = True
                for row in shape.table.rows:
                    html_table += "<tr>"
                    for cell in row.cells:
                        if first_row:
                            html_table += "<th>" + html.escape(cell.text) + "</th>"
                        else:
                            html_table += "<td>" + html.escape(cell.text) + "</td>"
                    html_table += "</tr>"
                    first_row = False
                html_table += "</table></body></html>"

                # Note: This would require a separate HTML to Markdown converter function
                # In this version, I'm assuming a convert_html_to_md function exists
                md_content += (
                    "\n" + convert_html_to_md(html_table).text_content.strip() + "\n"
                )

            # Text areas
            elif shape.has_text_frame:
                if shape == title:
                    md_content += "# " + shape.text.lstrip() + "\n"
                else:
                    md_content += shape.text + "\n"

        md_content = md_content.strip()
        if slide.has_notes_slide:
            md_content += "\n\n### Notes:\n"
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                md_content += notes_frame.text
            md_content = md_content.strip()

    return DocumentConverterResult(
        title=None,
        text_content=md_content.strip(),
    )


def _get_metadata(local_path):
    """
    Extract metadata from media files using exiftool

    Args:
        local_path: Path to the media file

    Returns:
        Dictionary of metadata if exiftool is available, None otherwise
    """
    exiftool = shutil.which("exiftool")
    if not exiftool:
        return None
    else:
        try:
            result = subprocess.run(
                [exiftool, "-json", local_path], capture_output=True, text=True
            ).stdout
            return json.loads(result)[0]
        except Exception:
            return None


def _transcribe_audio(local_path) -> str:
    """
    Transcribe audio using Google's speech recognition

    Args:
        local_path: Path to the audio file

    Returns:
        Transcription as a string
    """
    # Lazy import to avoid heavy deps / ffmpeg checks during program startup.
    import speech_recognition as sr

    recognizer = sr.Recognizer()
    with sr.AudioFile(local_path) as source:
        audio = recognizer.record(source)
        return recognizer.recognize_google(audio).strip()


def WavConverter(local_path) -> Union[None, dict]:
    """
    Converts WAV files to markdown via extraction of metadata (if `exiftool` is installed),
    and speech transcription (if `speech_recognition` is installed).

    Args:
        local_path: Path to the WAV file
        file_extension: Extension of the file (default: ".wav")

    Returns:
        None if not a WAV file, otherwise a dictionary with title and text_content
    """
    # Bail if not a WAV file

    md_content = ""

    # Add metadata
    metadata = _get_metadata(local_path)
    if metadata:
        for f in [
            "Title",
            "Artist",
            "Author",
            "Band",
            "Album",
            "Genre",
            "Track",
            "DateTimeOriginal",
            "CreateDate",
            "Duration",
        ]:
            if f in metadata:
                md_content += f"{f}: {metadata[f]}\n"

    # Transcribe
    try:
        transcript = _transcribe_audio(local_path)
        md_content += "\n\n### Audio Transcript:\n" + (
            "[No speech detected]" if transcript == "" else transcript
        )
    except Exception:
        md_content += (
            "\n\n### Audio Transcript:\nError. Could not transcribe this audio."
        )

    return DocumentConverterResult(
        title=None,
        text_content=md_content.strip(),
    )


def Mp3Converter(local_path: str, extension: str, **kwargs):
    # Lazy import to avoid pydub (and its ffmpeg detection) during program startup.
    import pydub

    md_content = ""

    # Add metadata
    metadata = _get_metadata(local_path)
    if metadata:
        for f in [
            "Title",
            "Artist",
            "Author",
            "Band",
            "Album",
            "Genre",
            "Track",
            "DateTimeOriginal",
            "CreateDate",
            "Duration",
        ]:
            if f in metadata:
                md_content += f"{f}: {metadata[f]}\n"

    # Transcribe
    handle, temp_path = tempfile.mkstemp(suffix=".wav")
    os.close(handle)
    try:
        if extension.lower() == ".mp3":
            sound = pydub.AudioSegment.from_mp3(local_path)
        else:
            sound = pydub.AudioSegment.from_file(local_path, format="m4a")
        sound.export(temp_path, format="wav")

        _args = dict()
        _args.update(kwargs)
        _args["file_extension"] = ".wav"

        try:
            transcript = super()._transcribe_audio(temp_path).strip()
            md_content += "\n\n### Audio Transcript:\n" + (
                "[No speech detected]" if transcript == "" else transcript
            )
        except Exception:
            md_content += (
                "\n\n### Audio Transcript:\nError. Could not transcribe this audio."
            )

    finally:
        os.unlink(temp_path)

    # Return the result
    return DocumentConverterResult(
        title=None,
        text_content=md_content.strip(),
    )


def ZipConverter(local_path: str, **kwargs):
    """
    Extracts ZIP files to a permanent local directory and returns a listing of extracted files.
    """
    md = MarkItDown(enable_plugins=False)
    md_content = md.convert(local_path).text_content

    return DocumentConverterResult(
        title="Extracted Files", text_content=md_content.strip()
    )
